from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import os
from models import enums
import requests
from io import BytesIO
from pptx.dml.color import RGBColor

def set_background(slide, color_hex):
    if color_hex:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex.lstrip("#"))

def set_text_style(text_frame, font_name, font_color):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if font_name:
                run.font.name = font_name
            if font_color:
                run.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))

def set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color):
    slide.shapes.title.text = slide_data.get("title", "")
    set_background(slide, background_color)
    set_text_style(slide.shapes.title.text_frame, font_name, font_color)

def build_pptx(presentation_id: int, slides: list, config: dict):
    prs = Presentation()
    
    font_name = config.get("font_name", "Arial")
    font_color = config.get("font_color", "#000000")
    background_color = config.get("background_color", "#FFFFFF")
    
    for slide_data in slides:
        layout = slide_data.get("layout", enums.SlideLayout.title.value)
        
        if layout == enums.SlideLayout.title.value:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
            set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color)
        
        elif layout == enums.SlideLayout.bullet.value:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Bullet slide
            set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color)
            body = slide.placeholders[1].text_frame
            for bullet in slide_data.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0        
        elif layout == enums.SlideLayout.two_column.value:
            slide = prs.slides.add_slide(prs.slide_layouts[3])  # Use default two-content layout
            set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color)

            # Get the content placeholders (usually index 1 and 2)
            left_placeholder = slide.placeholders[1]
            right_placeholder = slide.placeholders[2]

            # Clear existing content
            left_placeholder.text = ""
            right_placeholder.text = ""

            # Add left content
            for para in slide_data.get("left", "").split("\n"):
                p = left_placeholder.text_frame.add_paragraph()
                p.text = para.strip()
                p.level = 0

            # Add right content
            for para in slide_data.get("right", "").split("\n"):
                p = right_placeholder.text_frame.add_paragraph()
                p.text = para.strip()
                p.level = 0
        
        elif layout == enums.SlideLayout.image.value:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color)
            image_url = slide_data.get("image_url")
            
            if image_url:
                try:
                    # Download image
                    response = requests.get(image_url)
                    image_stream = BytesIO(response.content)
                    
                    # Add image to slide
                    slide.shapes.add_picture(
                        image_stream,
                        Inches(2), Inches(2),
                        width=Inches(6)
                    )
                except Exception as e:
                    print(f"Failed to add image: {str(e)}")
                    # Add placeholder rectangle if image fails
                    slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(2), Inches(2),
                        Inches(6), Inches(4)
                    )
            else:
                print(f"No image URL provided for slide: {slide_data.get('title')}")
                # Add placeholder rectangle for placeholder URLs
                slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(2), Inches(2),
                    Inches(6), Inches(4)
                )

    # Ensure output directory exists
    os.makedirs("output", exist_ok=True)
    path = os.path.abspath(f"./output/presentation_{presentation_id}.pptx")
    prs.save(path)
    return path