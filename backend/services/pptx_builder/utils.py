from pptx.dml.color import RGBColor

def set_background(slide, color_hex):
    if color_hex:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex.lstrip("#"))

def set_text_style(text_frame, font_name, font_color):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if font_name:
                run.font.name = font_name
            if font_color:
                run.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))

def set_slide_title_and_style(slide, slide_data, background_color, font_name, font_color):
    slide.shapes.title.text = slide_data.get("title", "")
    set_background(slide, background_color)
    set_text_style(slide.shapes.title.text_frame, font_name, font_color)