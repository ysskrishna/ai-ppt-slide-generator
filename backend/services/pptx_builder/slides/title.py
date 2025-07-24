from pptx.util import Inches
from services.pptx_builder.utils import set_slide_title_and_style

class TitleSlideStrategy:
    def add_slide(self, prs, slide_data, config, **kwargs):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_slide_title_and_style(slide, slide_data, config.get("background_color", "#FFFFFF"), config.get("font_name", "Arial"), config.get("font_color", "#000000"))
        footer = slide_data.get("footer_text") or kwargs.get("footer_text")
        if footer:
            left = top = width = height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, prs.slide_height - Inches(0.7), prs.slide_width - Inches(1), height)
            tf = txBox.text_frame
            tf.text = footer
        return slide 