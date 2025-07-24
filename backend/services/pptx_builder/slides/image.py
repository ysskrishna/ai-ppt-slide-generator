from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO
from services.pptx_builder.utils import set_slide_title_and_style

class ImageSlideStrategy:
    def add_slide(self, prs, slide_data, config, **kwargs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_title_and_style(slide, slide_data, config.get("background_color", "#FFFFFF"), config.get("font_name", "Arial"), config.get("font_color", "#000000"))
        image_url = slide_data.get("image_url")
        logger = kwargs.get("logger")
        if image_url:
            try:
                response = requests.get(image_url)
                image_stream = BytesIO(response.content)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(2), Inches(2),
                    width=Inches(6)
                )
            except Exception as e:
                if logger:
                    logger.error(f"Failed to add image: {str(e)}")
                else:
                    print(f"Failed to add image: {str(e)}")
                slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(2), Inches(2),
                    Inches(6), Inches(4)
                )
        else:
            if logger:
                logger.warning(f"No image URL provided for slide: {slide_data.get('title')}")
            else:
                print(f"No image URL provided for slide: {slide_data.get('title')}")
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(2), Inches(2),
                Inches(6), Inches(4)
            )
        return slide 